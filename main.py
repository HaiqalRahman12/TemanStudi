import shutil
import os
import tempfile
import fitz  # PyMuPDF
from pptx import Presentation
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List
from huggingface_hub import hf_hub_download

# Import Engine dari folder app/services
# (Struktur folder akan dibuat otomatis oleh Notebook Kaggle nanti)
from app.services.ai_engine import init_engine, get_engine

# --- KONFIGURASI MODEL ---
# Menggunakan Qwen 4B sesuai request
REPO_ID = "unsloth/Qwen3-4B-Instruct-2507-GGUF"
FILENAME = "Qwen3-4B-Instruct-2507-Q5_K_M.gguf"
# -------------------------

app = FastAPI(title="TemanStudi Kaggle Service")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class FlashcardItem(BaseModel):
    pertanyaan: str
    jawaban: str

class AIResponse(BaseModel):
    status: str
    pesan: str
    data: List[FlashcardItem]

# --- Helper Logic ---
def extract_from_pdf(path, s, e):
    doc = fitz.open(path)
    text = ""
    s_idx, e_idx = max(0, s-1), min(len(doc), e)
    for i in range(s_idx, e_idx):
        text += doc.load_page(i).get_text() + "\n"
    return text

def extract_from_pptx(path, s, e):
    prs = Presentation(path)
    text = ""
    s_idx, e_idx = max(0, s-1), min(len(prs.slides), e)
    for i in range(s_idx, e_idx):
        slide = prs.slides[i]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text.strip() + ". "
        text += "\n"
    return text

# --- Startup Event ---
@app.on_event("startup")
async def startup_event():
    print("⏳ [Main.py] Sedang mendownload model Qwen 4B... (Ini berjalan di background)")
    try:
        # Download ke folder ./models di working directory Kaggle
        os.makedirs("models", exist_ok=True)
        model_path = hf_hub_download(
            repo_id=REPO_ID,
            filename=FILENAME,
            local_dir="./models"
        )
        print(f"✅ [Main.py] Model Downloaded: {model_path}")
        
        # Init Engine
        init_engine(model_path)
        print("✅ [Main.py] AI Engine Siap & Loaded di GPU!")
    except Exception as e:
        print(f"❌ [Main.py] Gagal Init Engine: {e}")

# --- Endpoints ---
@app.post("/generate", response_model=AIResponse)
async def generate_endpoint(
    file: UploadFile = File(...),
    start_page: int = Form(1),
    end_page: int = Form(10)
):
    filename = file.filename.lower()
    if not (filename.endswith(".pdf") or filename.endswith(".pptx")):
        raise HTTPException(400, "Format harus PDF atau PPTX")

    ext = ".pdf" if filename.endswith(".pdf") else ".pptx"
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        shutil.copyfileobj(file.file, tmp)
        tmp_path = tmp.name

    try:
        full_text = ""
        if "pdf" in ext:
            full_text = extract_from_pdf(tmp_path, start_page, end_page)
        else:
            full_text = extract_from_pptx(tmp_path, start_page, end_page)

        if len(full_text) < 50:
            return AIResponse(status="error", pesan="Teks kosong", data=[])

        engine = get_engine()
        if engine is None:
             raise HTTPException(503, "AI Engine belum siap/sedang loading model. Coba 1 menit lagi.")
             
        results = engine.generate_flashcards(full_text)
        
        final_data = [FlashcardItem(pertanyaan=r['question'], jawaban=r['answer']) for r in results]
        return AIResponse(status="success", pesan="OK", data=final_data)

    finally:
        if os.path.exists(tmp_path): os.remove(tmp_path)

@app.get("/")
def health():
    engine = get_engine()
    status = "Ready" if engine else "Loading Model..."
    return {"status": status, "model": FILENAME}